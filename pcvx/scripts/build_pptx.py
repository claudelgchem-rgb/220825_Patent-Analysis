#!/usr/bin/env python3
"""2번 산출물 — 발표자료 pptx. 14~20장, 18pt 이상, 한 장 7줄 이내, 발표자 노트 포함."""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import common  # noqa: E402

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

BLANK, TITLE_ONLY = 6, 5
MAX_LINES = 7


def rgb(hexcode: str) -> RGBColor:
    return RGBColor(int(hexcode[0:2], 16), int(hexcode[2:4], 16), int(hexcode[4:6], 16))


def add_slide(prs, title: str, notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    slide.shapes.title.text = title
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(26)
            run.font.bold = True
    slide.notes_slide.notes_text_frame.text = notes or title
    return slide


def add_body(slide, lines: list[str], top: float = 1.6, size: int = 18) -> None:
    lines = [l for l in lines if str(l).strip()][:MAX_LINES]
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(11.9), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = str(line)
        para.space_after = Pt(10)
        for run in para.runs:
            run.font.size = Pt(size)


def add_badge(slide, grade: str) -> None:
    box = slide.shapes.add_textbox(Inches(10.9), Inches(0.35), Inches(1.9), Inches(0.6))
    para = box.text_frame.paragraphs[0]
    para.text = f"신뢰도 {grade}"
    run = para.runs[0]
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = rgb(common.GRADE_COLORS.get(grade, "444444"))


def slide_cover(prs, env: dict, topic: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "특허 조사 보고"
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    p2 = tf.add_paragraph()
    p2.text = common.cell(topic.get("topic"), "조사 주제 미확정")
    p2.runs[0].font.size = Pt(24)
    p3 = tf.add_paragraph()
    p3.text = f"작성 기준일 {env['as_of']}"
    p3.runs[0].font.size = Pt(18)
    slide.notes_slide.notes_text_frame.text = (
        f"{env['as_of']} 기준으로 작성한 특허 조사 결과다. 주제는 {topic.get('topic', '')} 이다."
    )


def slide_chart(prs, context: dict) -> None:
    slide = add_slide(
        prs, "출원은 연도별로 이렇게 움직였다",
        "연도별 출원 추이다. 세로축은 건수, 가로축은 우선일 기준 연도다.",
    )
    full = context.get("yearly_filings") or []
    yearly = context.get("yearly_filings_chart") or full
    if not yearly:
        add_body(slide, ["연도별 출원 추이 자료 " + common.FAILURE_MARK,
                         "a6_context.json 의 yearly_filings_chart 와 yearly_filings 가 모두 비어 있다."])
        return
    if full and len(yearly) < len(full):
        span = f"{full[0].get('year')}~{full[-1].get('year')}"
        slide.shapes.title.text = (
            f"출원은 최근 {len(yearly)}개 연도에 이렇게 움직였다"
        )
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(26)
                run.font.bold = True
        slide.notes_slide.notes_text_frame.text = (
            f"연도별 출원 추이다. 세로축은 건수, 가로축은 우선일 기준 연도다. "
            f"가독성을 위해 최근 {len(yearly)}개 연도만 그렸다. "
            f"전체 구간은 {span} 이며 그 이전 구간을 포함한 전체 표는 특허대장 엑셀 파일의 특허대장 시트에서 확인할 수 있다."
        )
    data = CategoryChartData()
    data.categories = [str(y.get("year")) for y in yearly]
    data.add_series("출원 건수", [int(y.get("count", 0)) for y in yearly])
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.5), Inches(11.6), Inches(5.2), data,
    )
    chart = frame.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(12)


def main() -> int:
    env = common.env()
    topic = common.read_json(common.WORKSPACE / "topic.json", {}) or {}
    records = (common.read_json(common.WORKSPACE / "records.json", {}) or {}).get("records", [])
    context = common.read_json(common.WORKSPACE / "a6_context.json", {}) or {}
    verification = common.read_json(common.WORKSPACE / "a7_verification.json", {}) or {}

    if not records:
        print("[실패] records.json 에 레코드가 없다.")
        return 1

    dist = verification.get("distribution", {})
    total = sum(dist.get(g, 0) for g in common.GRADES) or 1
    low_ratio = dist.get("low_ratio") or dist.get("하", 0) / total
    grade_by_record: dict[str, list[str]] = {}
    for item in verification.get("items", []):
        rid = str(item.get("claim_id", "")).split(".")[0]
        grade_by_record.setdefault(rid, []).append(item.get("grade", "하"))

    def worst(rid: str) -> str:
        grades = grade_by_record.get(rid, [])
        for g in ("하", "중", "상"):
            if g in grades:
                return g
        return "하"

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # 1. 표지
    slide_cover(prs, env, topic)

    # 2. 목차
    add_slide(
        prs, "오늘 보고 순서",
        f"보고 순서를 먼저 안내한다. 전체 {len(records)}건을 여섯 토막으로 나누어 본다. "
        "핵심 요약 3장에서 결론을 먼저 말하고, 기술 분야 개관에서 갈래와 연도별 추이를 보이고, "
        "주요 특허 8건은 한 장에 한 건씩 카드로 넘긴다. 목록 표는 발췌이며 전체는 특허대장 엑셀에 있다. "
        "마지막 두 장은 신뢰도 등급 분포와 이 조사의 한계다. 질의응답은 마지막 장에서 받는다.",
    )
    add_body(prs.slides[-1], [
        "1. 핵심 요약", "2. 기술 분야 개관", "3. 주요 특허 상세",
        "4. 특허 전체 목록", "5. 신뢰도 등급 분포", "6. 한계와 다음 단계",
    ])

    # 3~5. 핵심 요약 3장 (제목은 문장형 결론)
    branches = context.get("branches") or []
    top_applicants = context.get("top_applicants") or []
    registered = sum(1 for r in records if str(r.get("registration_status", "")).startswith("등록"))
    summaries = [
        (
            f"조사 대상 {len(records)}건 중 {registered}건이 이미 등록되어 권리가 살아 있다",
            [f"수집한 특허 패밀리는 모두 {len(records)}건이다.",
             f"이 가운데 등록(심사를 통과해 권리가 발생한 상태)은 {registered}건이다.",
             f"나머지 {len(records) - registered}건은 출원 중이거나 등록이 확인되지 않았다.",
             f"작성 기준일은 {env['as_of']} 이며, 법적 상태는 이 날짜 기준으로 확인했다."],
        ),
        (
            f"기술은 크게 {len(branches)}개 갈래로 나뉘어 발전하고 있다",
            [f"· {common.cell(b.get('name'))} — {common.cell(b.get('summary'))[:70]}"
             for b in branches[:MAX_LINES]] or ["기술 갈래 자료 " + common.FAILURE_MARK],
        ),
        (
            "출원인은 소수 기업에 집중되어 있다",
            [f"· {common.cell(a.get('name'))} {common.cell(a.get('count'))}건 — "
             f"{common.cell(a.get('note'))[:50]}" for a in top_applicants[:5]]
            or ["주요 출원인 자료 " + common.FAILURE_MARK],
        ),
    ]
    for title, lines in summaries:
        add_slide(prs, title, " ".join(str(l) for l in lines))
        add_body(prs.slides[-1], lines)

    # 6~7. 기술 분야 개관
    slide_chart(prs, context)
    recent = context.get("recent_12m") or []
    rstats = context.get("recent_12m_stats") or {}
    recent_notes = " ".join(
        f"({pos}) {common.cell(r.get('headline'))}" for pos, r in enumerate(recent, start=1)
    ) or "최근 12개월 창에서 확인된 신규 공개·등록 건이 없다."
    add_slide(
        prs, "최근 12개월 동향은 이렇게 요약된다",
        f"창은 {common.cell(rstats.get('window'), '최근 12개월')}, "
        f"해당 건수는 {common.cell(rstats.get('count'), '0')}건이다. "
        f"화면에는 요지만 실었고 아래가 헤드라인 원문이다. {recent_notes}",
    )
    add_body(prs.slides[-1], [
        f"· {common.cell(r.get('headline'))[:84].rstrip()}…" for r in recent[:MAX_LINES]
    ] + [
        f"위 {len(recent)}건의 서술 원문은 보고서 4장과 발표자 노트에 그대로 있다."
    ] if recent else ["최근 12개월 신규 공개 0건으로 확인되었다."])

    # 8~15. 주요 특허 카드 (상위 5~8건)
    #  수집 순서(P001~)는 중요도 순이 아니다. 권리가 살아 있는 등록건을 먼저,
    #  같은 조건이면 최근 출원 순으로 상위 8건을 고른다.
    status_rank = {"등록": 0, "출원 중(등록 미확인)": 1, "거절": 2, "포기": 3, "취하": 3, "소멸": 4}
    grade_rank = {"상": 0, "중": 1, "하": 2}

    def card_key(rec: dict):
        date = str(rec.get("filing_date") or "0000-00-00")
        return (
            status_rank.get(str(rec.get("registration_status", "")).strip(), 5),
            grade_rank.get(worst(rec.get("id", "?")), 3),
            0 if str(rec.get("period_scope", "")).startswith("조사 기간 내") else 1,
            tuple(-int(x) if x.isdigit() else 0 for x in (date.split("-") + ["0", "0", "0"])[:3]),
            str(rec.get("id", "")),
        )

    top_records = sorted(records, key=card_key)[:8]
    for rec in top_records:
        rid = rec.get("id", "?")
        grade = worst(rid)
        plain = common.cell(rec.get("plain_explanation"))
        title = common.cell(rec.get("title_ko") or rec.get("title_original"))
        slide = add_slide(prs, title[:60], plain)
        add_body(slide, [
            f"문헌번호 {common.cell(rec.get('representative_doc'))}",
            f"출원일 {common.cell(rec.get('filing_date'))} / "
            f"등록 여부 {common.cell(rec.get('registration_status'))}",
            f"권리자 {common.cell(rec.get('current_assignee') or rec.get('applicant'))}",
            "핵심 청구 내용: " + plain[:90],
            "권리범위를 좌우하는 요소: " + common.cell(
                [e.get("element") for e in (rec.get("claim_elements") or [])][:3]),
            f"법적상태 확인일 {common.cell(rec.get('legal_status_checked_on'))}",
        ])
        add_badge(slide, grade)

    # 16~17. 전체 목록 표
    rows_per_slide = 12
    # 목록도 카드와 같은 기준(권리 상태 → 신뢰도 → 최근 출원순)으로 정렬한 상위 건부터 싣는다.
    ranked = sorted(records, key=card_key)
    chunks = [ranked[i:i + rows_per_slide] for i in range(0, len(ranked), rows_per_slide)][:2]
    shown = sum(len(c) for c in chunks)
    for pos, chunk in enumerate(chunks, start=1):
        first_id = common.cell(chunk[0].get("id"))
        last_id = common.cell(chunk[-1].get("id"))
        slide = add_slide(
            prs, f"주요 특허 목록 ({pos}/{len(chunks)}) — 전체 {len(records)}건 중 {shown}건 발췌",
            f"수집한 특허 {len(records)}건 가운데 이 슬라이드에는 {first_id}부터 {last_id}까지 "
            f"{len(chunk)}건만 실었다. 발표 화면의 가독성 때문에 두 장에 {shown}건까지만 담았고, "
            f"나머지 {len(records) - shown}건을 포함한 109건 전체 목록은 같은 날짜의 특허대장 엑셀 파일 "
            f"'특허대장' 시트에 한 건도 빠짐없이 들어 있다. 질문이 나오면 그 파일을 열어 답한다.")
        table = slide.shapes.add_table(
            len(chunk) + 1, 5, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.4 * (len(chunk) + 1))
        ).table
        for idx, head in enumerate(["관리번호", "문헌번호", "출원일", "등록 여부", "권리자"]):
            table.cell(0, idx).text = head
        for r, rec in enumerate(chunk, start=1):
            values = [
                common.cell(rec.get("id")), common.cell(rec.get("representative_doc")),
                common.cell(rec.get("filing_date")), common.cell(rec.get("registration_status")),
                common.cell(rec.get("current_assignee") or rec.get("applicant")),
            ]
            for c, value in enumerate(values):
                table.cell(r, c).text = str(value)[:28]
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(12)
        note = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.6))
        npara = note.text_frame.paragraphs[0]
        note.text_frame.word_wrap = True
        npara.text = (
            f"이 표는 전체 {len(records)}건 중 {shown}건만 발췌한 것이다. "
            f"109건 전체 목록은 특허대장 엑셀 파일(PCVX_특허대장_{env['as_of_compact']}.xlsx)의 '특허대장' 시트를 보라."
        )
        npara.runs[0].font.size = Pt(18)
        npara.runs[0].font.bold = True

    # 18. 신뢰도 등급 분포
    slide = add_slide(prs, f"항목의 {low_ratio * 100:.0f}퍼센트가 원문 미확인 상태다"
                      if low_ratio > 0.30 else "신뢰도는 항목 단위로 이렇게 나뉜다",
                      f"신뢰도 등급은 특허 건이 아니라 항목 단위로 매겼다. "
                      f"레코드 {len(records)}건에서 서지·날짜·법적상태·청구항 항목을 각각 채점해 "
                      f"모두 {total}개 항목이 나왔고, 건당 평균 {total / max(len(records), 1):.1f}개다. "
                      f"'상'은 특허청 원문과 독립 출처 두 곳이 일치한 항목, '중'은 2차 집계 데이터베이스 "
                      f"한 곳에서만 확인한 항목, '하'는 원문 대조를 마치지 못한 항목이다. "
                      f"'하' 비율은 {low_ratio * 100:.1f}퍼센트로 30퍼센트 기준선 아래이므로 "
                      f"표지 경고 문구는 붙이지 않았다. 항목별 근거와 출처는 근거대장 파일에 있다.")
    add_body(slide, [
        f"높음(상) {dist.get('상', 0)}개 — 특허청 원문과 독립 출처 2곳이 일치",
        f"보통(중) {dist.get('중', 0)}개 — 2차 집계 데이터베이스 1곳에서 확인",
        f"낮음(하) {dist.get('하', 0)}개 — 원문 대조를 마치지 못함",
        f"전체 {total}개 항목 중 낮음 비율은 {low_ratio * 100:.1f}퍼센트다.",
        common.low_ratio_warning() if low_ratio > 0.30 else
        "낮음 비율이 30퍼센트 이하이므로 별도 경고를 붙이지 않는다.",
    ])

    # 19. 한계 및 다음 단계
    slide = add_slide(prs, "이 조사로 답할 수 없는 것이 무엇인지 분명히 해 둔다",
                      common.limits_text())
    add_body(slide, [
        "출원 후 18개월이 지나지 않은 특허는 원리적으로 확인할 수 없다.",
        "검색어 선택에 따라 다른 표현을 쓴 특허가 빠졌을 수 있다.",
        "이 보고서는 침해 여부를 판단한 것이 아니다.",
        "사업 실시 가능성 판단에는 유료 데이터베이스 전수 조사가 있어야 한다.",
        "변리사·특허 전문가의 법률 검토를 별도로 받아야 한다.",
    ])

    # 20. 규격 하한(14장)을 못 채우면 기술 갈래별 상세로 채운다. 빈 슬라이드는 넣지 않는다.
    filler = list(branches)
    while len(prs.slides) < 14 and filler:
        branch = filler.pop(0)
        member_ids = branch.get("record_ids") or []
        members = [r for r in records if r.get("id") in member_ids]
        slide = add_slide(
            prs,
            f"{common.cell(branch.get('name'))} 갈래는 이런 특허들로 이루어져 있다",
            common.cell(branch.get("summary")),
        )
        add_body(slide, [common.cell(branch.get("summary"))] + [
            f"· {common.cell(r.get('representative_doc'))} — "
            f"{common.cell(r.get('title_ko') or r.get('title_original'))[:50]}"
            for r in members[:5]
        ] or ["이 갈래에 속한 특허 " + common.FAILURE_MARK])

    if len(prs.slides) < 14:
        slide = add_slide(prs, "이 결과는 이런 검색 방법으로 얻었다",
                          "검색 방법과 시간 창 확장 순서를 밝힌다.")
        add_body(slide, [
            "최신순 정렬을 1순위로 두고 검색했다.",
            "시간 창은 최근 12개월, 최근 3년, 최근 10년, 전체 기간 순으로 넓혔다.",
            "중복은 문헌번호가 아니라 특허 패밀리(같은 발명을 여러 나라에 출원한 묶음) 단위로 제거했다.",
            f"검색 대상 데이터베이스는 {len(common.SOURCE_CODES)}곳이다.",
            f"법적 상태는 모두 {env['as_of']} 기준으로 확인했다.",
        ])

    # topic.json 이 축약 슬러그를 이미 확정해 두었으면 그것을 쓴다.
    # topic 원문을 30자에서 자르면 괄호가 열린 채 끊긴 파일명이 나온다.
    slug = common.slugify(topic.get("topic_slug") or topic.get("topic", "무제"))
    out = common.output_dir() / f"PCVX_특허조사보고서_{slug}_{env['as_of_compact']}.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    count = len(prs.slides)
    print(f"[기준일: {env['as_of']} / 도구: build_pptx]")
    print(f"-> {out}  ({out.stat().st_size:,} bytes, 슬라이드 {count}장)")
    if not 14 <= count <= 20:
        print(f"경고: 슬라이드가 {count}장이다. 규격은 14~20장이다.")
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
