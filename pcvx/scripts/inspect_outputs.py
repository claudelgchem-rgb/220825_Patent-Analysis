#!/usr/bin/env python3
"""산출물 실물 검사 — 파일을 실제로 열어 문단 수·슬라이드 수·행 수·빈 셀을 센다 (G7/G8).

  python3 pcvx/scripts/inspect_outputs.py
  python3 pcvx/scripts/inspect_outputs.py --strict   # 빈 항목·미완성 문장까지
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import common  # noqa: E402

# 청구항은 구성요소별로 줄바꿈하므로 줄 끝 쉼표는 정상이다. 조사로 끝나는 경우만 미완성으로 본다.
DANGLING = re.compile(r"(은|는|이|가|을|를|와|과|의|에|으로|로|및|또는)$")


def inspect_docx(path: Path, strict: bool) -> list[str]:
    from docx import Document

    problems: list[str] = []
    doc = Document(path)
    paragraphs = [p for p in doc.paragraphs if p.text.strip()]
    print(f"  {path.name}: 문단 {len(paragraphs)}개, 표 {len(doc.tables)}개")
    if not paragraphs:
        problems.append(f"- {path.name}: 본문 문단이 0개다.")
    if len(paragraphs) < 60:
        problems.append(f"- {path.name}: 본문 문단이 {len(paragraphs)}개다. 12~20쪽 분량에 크게 못 미친다.")

    text = "\n".join(p.text for p in doc.paragraphs)
    if common.limits_text().splitlines()[0][:40] not in text:
        problems.append(f"- {path.name}: 한계 문단(§7) 원문이 들어 있지 않다.")

    for t_idx, table in enumerate(doc.tables, start=1):
        for r_idx, row in enumerate(table.rows, start=1):
            for c_idx, cellobj in enumerate(row.cells, start=1):
                if not cellobj.text.strip():
                    problems.append(f"- {path.name}: 표 {t_idx} {r_idx}행 {c_idx}열이 빈 셀이다.")
    if strict:
        for p in paragraphs:
            body = p.text.strip()
            if len(body) > 25 and DANGLING.search(body):
                problems.append(f"- {path.name}: 미완성으로 끝나는 문장 — {body[-45:]!r}")
    return problems


def inspect_pptx(path: Path, strict: bool) -> list[str]:
    from pptx import Presentation

    problems: list[str] = []
    prs = Presentation(path)
    count = len(prs.slides)
    print(f"  {path.name}: 슬라이드 {count}장")
    if not 14 <= count <= 20:
        problems.append(f"- {path.name}: 슬라이드가 {count}장이다. 규격은 14~20장이다.")
    for idx, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        has_table = any(shape.has_table for shape in slide.shapes)
        has_chart = any(shape.has_chart for shape in slide.shapes)
        if not texts and not has_table and not has_chart:
            problems.append(f"- {path.name}: {idx}번 슬라이드가 비어 있다.")
        body_lines = sum(len(t.splitlines()) for t in texts[1:])
        if body_lines > 7:
            problems.append(f"- {path.name}: {idx}번 슬라이드 본문이 {body_lines}줄이다. 7줄 이내로 줄인다.")
        if not slide.notes_slide.notes_text_frame.text.strip():
            problems.append(f"- {path.name}: {idx}번 슬라이드에 발표자 노트가 없다.")
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for r, row in enumerate(shape.table.rows, start=1):
                for c, cellobj in enumerate(row.cells, start=1):
                    if not cellobj.text.strip():
                        problems.append(f"- {path.name}: {idx}번 슬라이드 표 {r}행 {c}열이 빈 셀이다.")
        if strict:
            from pptx.util import Pt

            for shape in slide.shapes:
                if not shape.has_text_frame or shape.has_table:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None and run.font.size < Pt(18):
                            problems.append(
                                f"- {path.name}: {idx}번 슬라이드 글꼴이 "
                                f"{run.font.size.pt:.0f}pt 다. 18pt 이상이어야 한다 — {run.text[:30]!r}"
                            )
    return problems


def inspect_xlsx(path: Path, strict: bool) -> list[str]:
    from openpyxl import load_workbook

    problems: list[str] = []
    wb = load_workbook(path)
    print(f"  {path.name}: 시트 {len(wb.sheetnames)}개 — {', '.join(wb.sheetnames)}")
    for ws in wb.worksheets:
        rows = ws.max_row - 1
        if rows < 1:
            problems.append(f"- {path.name}: '{ws.title}' 시트에 데이터 행이 없다.")
            continue
        print(f"    {ws.title}: 데이터 {rows}행 x {ws.max_column}열")
        for row in ws.iter_rows(min_row=2):
            for cellobj in row:
                if cellobj.value is None or str(cellobj.value).strip() == "":
                    problems.append(
                        f"- {path.name}: '{ws.title}' {cellobj.coordinate} 가 빈 셀이다. "
                        "'[확보 실패]' + 사유를 쓴다."
                    )
    return problems


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--strict", action="store_true")
    args = ap.parse_args()

    out = common.output_dir()
    print(f"[기준일: {common.today()} / 검사: inspect_outputs] 대상 {out}")
    problems: list[str] = []
    found = {"docx": False, "pptx": False, "xlsx": False, "md": 0}

    for path in sorted(out.glob("PCVX_*")):
        suffix = path.suffix.lower().lstrip(".")
        if suffix == "docx":
            found["docx"] = True
            problems += inspect_docx(path, args.strict)
        elif suffix == "pptx":
            found["pptx"] = True
            problems += inspect_pptx(path, args.strict)
        elif suffix == "xlsx":
            found["xlsx"] = True
            problems += inspect_xlsx(path, args.strict)
        elif suffix == "md":
            found["md"] += 1
            text = path.read_text(encoding="utf-8")
            print(f"  {path.name}: {len(text.splitlines())}줄")
            if not text.strip():
                problems.append(f"- {path.name}: 내용이 비어 있다.")
            if not text.lstrip().startswith("[기준일:"):
                problems.append(f"- {path.name}: 첫 줄에 '[기준일: YYYY-MM-DD]' 표기가 없다.")

    for key in ("docx", "pptx", "xlsx"):
        if not found[key]:
            problems.append(f"- {key} 파일이 출력 디렉터리에 없다.")
    if args.strict and found["md"] < 2:
        problems.append(f"- 근거 대장과 감사 로그 두 개가 필요한데 md 파일이 {found['md']}개다.")

    if problems:
        print(f"\n위반 {len(problems)}건")
        for line in problems:
            print(f"  {line}")
        return 1
    print("통과: 위반 0건")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
